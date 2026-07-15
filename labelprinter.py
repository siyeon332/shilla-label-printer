import io
import os
import urllib.request
import streamlit as st
import barcode
from barcode.writer import ImageWriter

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# PDF 생성을 위한 reportlab 라이브러리
from reportlab.lib.pagesizes import letter, landscape
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.graphics.barcode import code128
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont

# --- [인터넷에서 굵은 한글 폰트 강제 다운로드 시스템] ---
# Streamlit 서버에 한글 폰트가 없으므로, 눈에 확 띄는 구글의 'Noto Sans KR Bold' (본고딕 굵은체)를 실시간 다운로드합니다.
FONT_NAME = "NotoSansKR-Bold"
FONT_FILE = "NotoSansKR-Bold.ttf"

@st.cache_resource
def load_korean_font():
    """웹 서버 환경에서도 완벽한 한글 출력을 위해 웹에서 굵은 고딕 폰트를 다운로드하여 등록합니다."""
    try:
        if not os.path.exists(FONT_FILE):
            # 구글 웹 폰트 저장소에서 진짜 두껍고 듬직한 본고딕 Bold 폰트 다운로드
            font_url = "https://github.com/google/fonts/raw/main/ofl/notosanskr/NotoSansKR%5Bwght%5D.ttf"
            urllib.request.urlretrieve(font_url, FONT_FILE)
        
        pdfmetrics.registerFont(TTFont(FONT_NAME, FONT_FILE))
        return True
    except Exception as e:
        return False

# 폰트 로드 실행
FONT_LOADED = load_korean_font()

# 만약 다운로드 실패 시 쓸 예비용 폰트 이름 설정
if not FONT_LOADED:
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    pdfmetrics.registerFont(UnicodeCIDFont("HYGothic-Medium"))
    FONT_NAME = "HYGothic-Medium"


# 페이지 기본 설정
st.set_page_config(
    page_title="신라면세점 라벨 원스톱 생성기",
    page_icon="🏷️",
    layout="centered"
)

# 세션 상태 초기화
if "barcode_bytes" not in st.session_state:
    st.session_state.barcode_bytes = None
if "barcode_val" not in st.session_state:
    st.session_state.barcode_val = ""

# --- [기능 1] 내부 메모리에서 바코드(Code 128) 이미지 파일 생성 (PPTX용) ---
def generate_code128_image(value: str) -> io.BytesIO:
    Code128 = barcode.get_barcode_class("code128")
    fp = io.BytesIO()
    barcode_image = Code128(value, writer=ImageWriter())
    barcode_image.write(
        fp,
        options={
            "module_width": 0.5,
            "module_height": 24,
            "font_size": 14,
            "text_distance": 10,
            "quiet_zone": 6,
            "dpi": 300,
            "write_text": True,
        }
    )
    fp.seek(0)
    return fp

# --- [기능 2] PPTX 라벨 일괄 생성 ---
def create_shilla_pptx(brand_name: str, label_type: str, total_qty: int, barcode_bytes: io.BytesIO) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(11.69)   # A4 가로
    prs.slide_height = Inches(8.27)
    blank_layout = prs.slide_layouts[6]

    import tempfile
    import os
    with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
        tmp.write(barcode_bytes.getvalue())
        tmp_barcode_path = tmp.name

    try:
        for i in range(1, total_qty + 1):
            slide = prs.slides.add_slide(blank_layout)
            
            # 1. 브랜드명 (크기 96, 맑은 고딕, 굵게 적용)
            brand_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10.69), Inches(1.8))
            tf_brand = brand_box.text_frame
            tf_brand.word_wrap = True
            p_brand = tf_brand.paragraphs[0]
            p_brand.text = brand_name
            p_brand.alignment = PP_ALIGN.CENTER
            p_brand.font.name = "맑은 고딕"
            p_brand.font.size = Pt(96)  # 요청하신 96 크기 반영!
            p_brand.font.bold = True    # 엄청 두껍고 힘 있게!
            
            # 2. PLT/BOX 번호 (크기 72, 맑은 고딕, 굵게 적용)
            plt_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.3), Inches(10.69), Inches(1.5))
            tf_plt = plt_box.text_frame
            tf_plt.word_wrap = True
            p_plt = tf_plt.paragraphs[0]
            p_plt.text = f"{label_type} NO. {total_qty}-{i}"
            p_plt.alignment = PP_ALIGN.CENTER
            p_plt.font.name = "맑은 고딕"
            p_plt.font.size = Pt(72)    # 96 크기에 맞춰 밸런스 있게 키웠습니다!
            p_plt.font.bold = True      # 굵게 적용!
            
            # 3. 바코드 이미지 위치 조정 (상단 글자가 커졌으므로 위치를 아래로 살짝 내림)
            img_width = Inches(8.5)
            img_height = Inches(3.8)
            img_left = (prs.slide_width - img_width) / 2
            img_top = Inches(4.0)
            slide.shapes.add_picture(tmp_barcode_path, img_left, img_top, width=img_width, height=img_height)
            
    finally:
        if os.path.exists(tmp_barcode_path):
            os.unlink(tmp_barcode_path)

    out_pptx = io.BytesIO()
    prs.save(out_pptx)
    out_pptx.seek(0)
    return out_pptx

# --- [기능 3] PDF 라벨 일괄 생성 ---
def create_shilla_pdf(brand_name: str, label_type: str, total_qty: int, barcode_value: str) -> io.BytesIO:
    pdf_buffer = io.BytesIO()
    pagesize = (841.89, 595.27) # A4 가로 크기
    c = canvas.Canvas(pdf_buffer, pagesize=pagesize)
    
    for i in range(1, total_qty + 1):
        # 1. 브랜드명 작성 (인터넷에서 받은 진짜 두껍고 든든한 고딕 폰트 사용!)
        c.setFont(FONT_NAME, 90)
        c.drawCentredString(pagesize[0]/2.0, 460, brand_name)
        
        # 2. PLT / BOX 번호 작성 (동일한 두껍고 선명한 폰트 적용)
        c.setFont(FONT_NAME, 68)
        c.drawCentredString(pagesize[0]/2.0, 360, f"{label_type} NO. {total_qty}-{i}")
        
        # 3. Code 128 바코드 그리기
        barcode_obj = code128.Code128(barcode_value, barWidth=1.8, barHeight=140, humanReadable=True)
        barcode_obj.fontName = "Helvetica-Bold" # 바코드 아래 숫자는 굵은 헬베티카로 선명하게!
        barcode_obj.fontSize = 16
        
        barcode_width = barcode_obj.width
        x_pos = (pagesize[0] - barcode_width) / 2.0
        y_pos = 90
        
        barcode_obj.drawOn(c, x_pos, y_pos)
        
        if i < total_qty:
            c.showPage()
            
    c.save()
    pdf_buffer.seek(0)
    return pdf_buffer


# --- 🖥️ 초간편 웹 화면 (Streamlit UI) ---
st.title("🏷️ 신라면세점 부착라벨 원스톱 생성기")
st.markdown("입력하신 수량과 동일하게 파일 페이지가 일괄 구성됩니다.")
st.write("---")

col1, col2 = st.columns(2)

with col1:
    label_type = st.radio("1. 라벨 형태 선택", ["PLT", "BOX"], horizontal=True)
    brand_name = st.text_input("2. 브랜드명 입력", value="아비브(ABIB)")

with col2:
    total_qty = st.number_input(f"3. 총 {label_type} 수량 (N)", min_value=1, value=5, step=1)
    barcode_val = st.text_input("4. 바코드 번호 입력", placeholder="예: 7851260079")

st.write(" ")
st.write(" ")

btn_generate_all = st.button("🚀 신라면세점 규격 라벨 파일 생성", type="primary", use_container_width=True)

if btn_generate_all:
    if not barcode_val.strip():
        st.error("바코드 번호를 입력해 주세요!")
    else:
        with st.spinner("라벨 내부 디자인 및 바코드 생성 처리 중..."):
            try:
                # 1. 공통 바코드 생성
                barcode_bytes = generate_code128_image(barcode_val.strip())
                st.session_state.barcode_bytes = barcode_bytes
                st.session_state.barcode_val = barcode_val.strip()
                
                # 2. PPTX 파일 생성
                pptx_output = create_shilla_pptx(
                    brand_name=brand_name.strip(),
                    label_type=label_type,
                    total_qty=total_qty,
                    barcode_bytes=barcode_bytes
                )
                
                # 3. PDF 파일 생성
                pdf_output = create_shilla_pdf(
                    brand_name=brand_name.strip(),
                    label_type=label_type,
                    total_qty=total_qty,
                    barcode_value=barcode_val.strip()
                )
                
                st.success("🎉 라벨 파일들이 성공적으로 디자인되었습니다!")
                
                st.image(barcode_bytes, caption=f"자동 생성된 Code 128 바코드 ({barcode_val})", width=320)
                
                st.write("---")
                st.markdown("### 📥 원하시는 포맷의 버튼을 눌러 다운로드하세요")
                
                dl_col1, dl_col2 = st.columns(2)
                
                with dl_col1:
                    st.download_button(
                        label="📥 PPTX 파일 다운로드",
                        data=pptx_output,
                        file_name=f"{brand_name}_라벨_{label_type}{total_qty}.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
                    
                with dl_col2:
                    st.download_button(
                        label="📥 PDF 파일 다운로드",
                        data=pdf_output,
                        file_name=f"{brand_name}_라벨_{label_type}{total_qty}.pdf",
                        mime="application/pdf",
                        use_container_width=True
                    )
                
            except Exception as e:
                st.error(f"라벨 자동 제작 실패: {e}")
